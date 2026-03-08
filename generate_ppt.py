"""
全球舆情态势感知平台 - 讲解PPT生成脚本（含截图占位区域）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ========== 颜色常量 ==========
DARK_BG = RGBColor(0x0B, 0x14, 0x26)        # 深蓝背景
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xFF)    # 科技蓝
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)    # 青绿色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xC4, 0xDE)     # 浅蓝灰
GOLD = RGBColor(0xFF, 0xD7, 0x00)           # 金色强调
RED_ALERT = RGBColor(0xFF, 0x44, 0x44)      # 警告红
CARD_BG = RGBColor(0x10, 0x1D, 0x35)        # 卡片背景
SECTION_BLUE = RGBColor(0x14, 0x28, 0x50)   # 章节页背景
PLACEHOLDER_BG = RGBColor(0x0D, 0x18, 0x2E) # 截图占位背景
PLACEHOLDER_BORDER = RGBColor(0x2A, 0x4A, 0x7A)  # 截图占位边框

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """设置幻灯片背景颜色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_decorated_card(slide, left, top, width, height, accent_color=ACCENT_BLUE):
    """添加带顶部装饰线的卡片"""
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color=CARD_BG)
    add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top, width - Inches(0.1), Inches(0.04), fill_color=accent_color)
    return card


def add_screenshot_placeholder(slide, left, top, width, height, label):
    """
    添加截图占位区域：虚线边框 + 居中标签文字
    使用后替换为实际截图即可
    """
    # 占位背景
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                    fill_color=PLACEHOLDER_BG, line_color=PLACEHOLDER_BORDER, line_width=Pt(2))
    # 对角线装饰（用两条细线模拟，表示这是占位区）
    # 左上到右下
    add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + height / 2 - Inches(0.005),
              width - Inches(0.3), Inches(0.01), fill_color=PLACEHOLDER_BORDER)
    # 居中标签
    add_textbox(slide, left, top + height / 2 - Inches(0.45), width, Inches(0.4),
                "[ 截图占位 ]", font_size=11, color=PLACEHOLDER_BORDER, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + height / 2 - Inches(0.1), width, Inches(0.5),
                label, font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    return box


def add_bottom_bar(slide):
    """添加底部装饰条"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_HEIGHT - Inches(0.06), SLIDE_WIDTH, Inches(0.06), fill_color=ACCENT_BLUE)


def add_page_number(slide, num, total):
    """添加页码"""
    add_textbox(slide, SLIDE_WIDTH - Inches(1.2), SLIDE_HEIGHT - Inches(0.5), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


def create_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    total_slides = 12

    # ====================================================
    # 第1页：封面
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.05), fill_color=ACCENT_BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.8), Inches(0.06), Inches(3.5), fill_color=ACCENT_BLUE)

    add_textbox(slide, Inches(2.0), Inches(1.8), Inches(9), Inches(1.2),
                "全球舆情态势感知平台", font_size=48, color=WHITE, bold=True)
    add_textbox(slide, Inches(2.0), Inches(3.0), Inches(9), Inches(0.8),
                "Global Public Opinion Situational Awareness Platform", font_size=20, color=ACCENT_CYAN)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(3.9), Inches(3), Inches(0.03), fill_color=ACCENT_BLUE)
    add_textbox(slide, Inches(2.0), Inches(4.2), Inches(8), Inches(0.5),
                "皇岗边检站  |  数智赋能 · 智慧管控", font_size=22, color=LIGHT_GRAY)

    add_bottom_bar(slide)

    # ====================================================
    # 第2页：背景与初衷
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "建设背景与初衷", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), fill_color=ACCENT_BLUE)

    # 左侧 - 战略部署
    add_decorated_card(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2), ACCENT_CYAN)
    add_textbox(slide, Inches(1.1), Inches(1.85), Inches(5.3), Inches(0.5),
                "⬡  站党委战略部署", font_size=22, color=ACCENT_CYAN, bold=True)

    items = [
        '2026年 — "十五五"开局之年、新皇岗口岸建成启用之年',
        "深圳举办 APEC 会议 — 核心安保之年",
        '打造"数智赋能的智慧新标杆"等五个示范新标杆',
        '"大站要有大作为" — 全面建设高质量发展强站',
    ]
    y = Inches(2.55)
    for item in items:
        add_textbox(slide, Inches(1.3), y, Inches(5.1), Inches(0.55),
                    f"▸  {item}", font_size=15, color=LIGHT_GRAY)
        y += Inches(0.65)

    # 右侧 - 痛点
    add_decorated_card(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.2), RED_ALERT)
    add_textbox(slide, Inches(7.2), Inches(1.85), Inches(5.3), Inches(0.5),
                "⚠  一线痛点", font_size=22, color=RED_ALERT, bold=True)

    pain_items = [
        '数据研判民警每天在境外网站"大海捞针"',
        "信息搜集全凭个人经验，无法传承复用",
        "换人就得重头来过，效率极低",
    ]
    y = Inches(2.55)
    for item in pain_items:
        add_textbox(slide, Inches(7.4), y, Inches(5.1), Inches(0.55),
                    f"▸  {item}", font_size=15, color=LIGHT_GRAY)
        y += Inches(0.65)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(4.6), Inches(5.1), Inches(0.03), fill_color=ACCENT_BLUE)
    add_textbox(slide, Inches(7.2), Inches(4.85), Inches(5.1), Inches(1.5),
                '▶  将先进经验"流程化"，用系统平台继承\n▶  让机器干繁琐搜集，解放警力做研判',
                font_size=15, color=ACCENT_CYAN)

    add_bottom_bar(slide)
    add_page_number(slide, 2, total_slides)

    # ====================================================
    # 第3页：安全可靠性
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "一、自主可控 · 安全可靠", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), fill_color=ACCENT_BLUE)

    # 左侧 - 技术可行
    add_decorated_card(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.5), ACCENT_CYAN)
    add_textbox(slide, Inches(1.1), Inches(1.85), Inches(5.3), Inches(0.5),
                "✓  技术上绝对可行", font_size=24, color=ACCENT_CYAN, bold=True)

    tech_items = [
        "完全自主研发，代码在我们自己手里",
        "可随时根据站里新需求修改升级",
        "无需依赖外部公司，灵活可控",
    ]
    y = Inches(2.6)
    for item in tech_items:
        add_textbox(slide, Inches(1.3), y, Inches(5.1), Inches(0.5),
                    f"▸  {item}", font_size=16, color=LIGHT_GRAY)
        y += Inches(0.6)

    # 右侧 - 安全可靠
    add_decorated_card(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(4.5), GOLD)
    add_textbox(slide, Inches(7.2), Inches(1.85), Inches(5.3), Inches(0.5),
                "🔒  安全上绝对可靠", font_size=24, color=GOLD, bold=True)

    safe_items = [
        '"只进不出"原则 — 只从外网公开渠道获取信息',
        "绝不会把内部数据传出去",
        "系统和数据全部署在国内自有服务器",
        "完全符合国家数据安全合规要求",
    ]
    y = Inches(2.6)
    for item in safe_items:
        add_textbox(slide, Inches(7.4), y, Inches(5.1), Inches(0.5),
                    f"▸  {item}", font_size=15, color=LIGHT_GRAY)
        y += Inches(0.55)

    # 底部核心原则
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.7),
              fill_color=SECTION_BLUE, line_color=ACCENT_BLUE, line_width=Pt(1.5))
    add_textbox(slide, Inches(3.5), Inches(6.35), Inches(6.3), Inches(0.6),
                "核心原则：只收信、不寄信 — 外部攻不进来，内部带不走",
                font_size=17, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide)
    add_page_number(slide, 3, total_slides)

    # ====================================================
    # 第4页：四大核心功能概览
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "二、四大核心功能", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), fill_color=ACCENT_BLUE)

    cards = [
        {
            "title": '"侦察兵" — 多源实时采集',
            "icon_color": ACCENT_BLUE,
            "items": ["全天候自动抓取全球30+主流媒体", "突破时差限制，24小时不间断", "上班即有最新涉边检资讯"],
            "pos": (Inches(0.8), Inches(1.6)),
        },
        {
            "title": '"翻译官" — 随身翻译 + AI 分析',
            "icon_color": ACCENT_CYAN,
            "items": ["自动翻译多语种新闻内容", "AI 一键分析热点重点", "智能生成舆情总结，一目了然"],
            "pos": (Inches(6.9), Inches(1.6)),
        },
        {
            "title": '"报警器" — 智能风险预警',
            "icon_color": RED_ALERT,
            "items": ["自定义敏感关键词监控", "高 / 中 / 低 三级风险分类", "突发事件第一时间红牌告警"],
            "pos": (Inches(0.8), Inches(4.2)),
        },
        {
            "title": '"作战地图" — 可视化态势',
            "icon_color": GOLD,
            "items": ["全球地图实时标注事件位置", "类似指挥中心监控大屏", "全球涉边舆情一目了然"],
            "pos": (Inches(6.9), Inches(4.2)),
        },
    ]

    for card in cards:
        x, y = card["pos"]
        add_decorated_card(slide, x, y, Inches(5.8), Inches(2.3), card["icon_color"])
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.5),
                    card["title"], font_size=20, color=card["icon_color"], bold=True)
        item_y = y + Inches(0.75)
        for item in card["items"]:
            add_textbox(slide, x + Inches(0.5), item_y, Inches(5.0), Inches(0.4),
                        f"▸  {item}", font_size=14, color=LIGHT_GRAY)
            item_y += Inches(0.42)

    add_bottom_bar(slide)
    add_page_number(slide, 4, total_slides)

    # ====================================================
    # 第5页：主仪表盘展示（截图页）
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "平台主界面总览", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), fill_color=ACCENT_BLUE)

    # 大截图占位 - 主仪表盘全貌
    add_screenshot_placeholder(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3),
                               "截取：主仪表盘全貌（首页 / ）\n包含左侧数据概览+文章列表、中央地图、右侧风控告警+趋势图")

    add_textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
                "▸ 左面板：数据概览 + 最新文章列表　　▸ 中央：全球态势地图　　▸ 右面板：风控告警 + 趋势分析",
                font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide)
    add_page_number(slide, 5, total_slides)

    # ====================================================
    # 第6页：多源采集详解 + 截图
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "功能详解 ①  多源实时采集", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), fill_color=ACCENT_BLUE)

    # 左侧 - 数据源分类文字
    source_groups = [
        ("港台媒体", ACCENT_BLUE, "联合新闻网、中时新闻网\n自由时报、ETtoday、TVBS …"),
        ("亚洲华语媒体", ACCENT_CYAN, "联合早报、星洲日报\n光华日报、中央社 …"),
        ("国际主流媒体", GOLD, "BBC中文、VOA中文\nCNN、Reuters、NHK …"),
        ("政府移民机构", RGBColor(0xAA, 0x77, 0xFF), "各国移民局官网\n政府公报、政策公告 …"),
    ]

    y_start = Inches(1.5)
    for i, (name, color, sources) in enumerate(source_groups):
        y = y_start + i * Inches(1.35)
        add_decorated_card(slide, Inches(0.8), y, Inches(4.5), Inches(1.15), color)
        add_textbox(slide, Inches(1.1), y + Inches(0.15), Inches(1.8), Inches(0.4),
                    name, font_size=16, color=color, bold=True)
        add_textbox(slide, Inches(2.9), y + Inches(0.15), Inches(2.2), Inches(0.85),
                    sources, font_size=11, color=LIGHT_GRAY)

    # 右侧 - 截图占位：最新文章列表 / 新闻源管理
    add_screenshot_placeholder(slide, Inches(5.6), Inches(1.5), Inches(7.0), Inches(3.0),
                               "截取：左面板 — 最新获取文章列表\n展示自动抓取的新闻标题、来源、时间")

    add_screenshot_placeholder(slide, Inches(5.6), Inches(4.7), Inches(7.0), Inches(2.3),
                               "截取：插件管理页面（设置 → 插件管理）\n展示新闻源站点开关、抓取方式配置")

    add_bottom_bar(slide)
    add_page_number(slide, 6, total_slides)

    # ====================================================
    # 第7页：AI分析 & 智能预警 + 截图
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "功能详解 ②  AI 分析 & 智能预警", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(3.5), Inches(0.04), fill_color=ACCENT_BLUE)

    # 上排左 - AI能力文字
    add_decorated_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.6), ACCENT_CYAN)
    add_textbox(slide, Inches(1.1), Inches(1.72), Inches(5.0), Inches(0.5),
                "AI 智能分析", font_size=20, color=ACCENT_CYAN, bold=True)

    ai_items = [
        "多语种新闻自动翻译为中文",
        "支持多种 LLM 提供商按需切换",
        "AI 一键生成舆情日报 / 周报",
        "自定义分析提示词，灵活适配",
    ]
    y = Inches(2.3)
    for item in ai_items:
        add_textbox(slide, Inches(1.3), y, Inches(4.8), Inches(0.4),
                    f"▸  {item}", font_size=14, color=LIGHT_GRAY)
        y += Inches(0.45)

    # 上排右 - AI总结截图
    add_screenshot_placeholder(slide, Inches(6.6), Inches(1.5), Inches(6.0), Inches(2.6),
                               "截取：AI 舆情总结页面\n展示 AI 生成的舆情分析报告内容")

    # 下排左 - 三级预警文字
    add_decorated_card(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(2.5), RED_ALERT)
    add_textbox(slide, Inches(1.1), Inches(4.62), Inches(5.0), Inches(0.5),
                "三级风险预警体系", font_size=20, color=RED_ALERT, bold=True)

    levels = [
        ("高风险", RGBColor(0xFF, 0x22, 0x22), "涉恐、涉暴等紧急事件"),
        ("中风险", RGBColor(0xFF, 0xAA, 0x00), "政策变动、群体事件"),
        ("低风险", RGBColor(0xFF, 0xDD, 0x44), "舆论趋势、一般报道"),
    ]
    y = Inches(5.2)
    for name, color, desc in levels:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), y, Inches(1.1), Inches(0.35), fill_color=color)
        add_textbox(slide, Inches(1.3), y + Inches(0.01), Inches(1.1), Inches(0.33),
                    name, font_size=12, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.6), y + Inches(0.01), Inches(3.4), Inches(0.35),
                    desc, font_size=13, color=LIGHT_GRAY)
        y += Inches(0.47)

    # 下排右 - 风控告警截图
    add_screenshot_placeholder(slide, Inches(6.6), Inches(4.4), Inches(6.0), Inches(2.5),
                               "截取：右面板 — 风控预警告警列表\n展示高/中/低风险告警条目和红牌提示")

    add_bottom_bar(slide)
    add_page_number(slide, 7, total_slides)

    # ====================================================
    # 第8页：可视化地图 + 新闻预览截图
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "功能详解 ③  作战地图 & 新闻预览", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(3.5), Inches(0.04), fill_color=ACCENT_BLUE)

    # 左侧 - 地图截图
    add_screenshot_placeholder(slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(4.0),
                               "截取：中央全球态势地图\n展示世界地图上的新闻事件标注点分布")

    add_textbox(slide, Inches(0.8), Inches(5.6), Inches(6.0), Inches(0.4),
                "▸ 全球事件实时标注，点击可查看详情",
                font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # 右侧 - 新闻预览截图
    add_screenshot_placeholder(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(4.0),
                               "截取：新闻预览弹窗（点击文章标题后）\n展示正文提取、翻译、截图等预览功能")

    add_textbox(slide, Inches(7.1), Inches(5.6), Inches(5.5), Inches(0.4),
                "▸ 智能正文提取 + 一键翻译 + 页面截图",
                font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # 底部功能说明
    add_decorated_card(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.9), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), Inches(6.2), Inches(11.2), Inches(0.7),
                "智能回退链：正文提取 → 页面截图 → 缓存读取　|　支持多种地图源切换（高德/OpenStreetMap/卫星图）　|　图片服务端代理中转",
                font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide)
    add_page_number(slide, 8, total_slides)

    # ====================================================
    # 第9页：实战案例 - APEC（截图页）
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "实战案例：APEC 预警", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), fill_color=GOLD)

    # 流程箭头
    steps = [
        ("设置关键词", 'APEC 会议在即\n设置"apec"为\n风险提示词', ACCENT_BLUE),
        ("自动抓取", '系统自动抓取\n"美国APEC组织\n高官访台"新闻', ACCENT_CYAN),
        ("预警触发", '民警上班即看到\n高风险告警\n第一时间掌握', RED_ALERT),
        ("成果产出", '立刻撰写分析文章\n文章已报送\n移民局', GOLD),
    ]

    card_w = Inches(2.65)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (title, desc, color) in enumerate(steps):
        x = x_start + i * (card_w + gap)
        add_decorated_card(slide, x, Inches(1.5), card_w, Inches(2.0), color)
        add_textbox(slide, x + Inches(0.15), Inches(1.72), card_w - Inches(0.3), Inches(0.4),
                    f"Step {i+1}  {title}", font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(2.2), card_w - Inches(0.3), Inches(1.1),
                    desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow_x = x + card_w + Inches(0.05)
            add_textbox(slide, arrow_x, Inches(2.2), gap - Inches(0.1), Inches(0.5),
                        "→", font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # 截图占位 - 展示实际的告警触发
    add_screenshot_placeholder(slide, Inches(0.8), Inches(3.8), Inches(5.8), Inches(3.0),
                               '截取：风控关键词设置页面\n展示"apec"等关键词的配置界面')

    add_screenshot_placeholder(slide, Inches(6.9), Inches(3.8), Inches(5.8), Inches(3.0),
                               '截取：APEC 相关告警/新闻\n展示系统抓取到的 APEC 相关新闻条目')

    add_bottom_bar(slide)
    add_page_number(slide, 9, total_slides)

    # ====================================================
    # 第10页：系统技术架构
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "系统技术架构", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), fill_color=ACCENT_BLUE)

    # 数据流
    flow_steps = [
        ("全球新闻源\n30+ 媒体", ACCENT_BLUE),
        ("智能采集引擎\n4种采集方式", ACCENT_CYAN),
        ("AI 处理层\n翻译·分析·预警", GOLD),
        ("MongoDB\n数据存储", RGBColor(0x55, 0xBB, 0x55)),
        ("可视化仪表盘\n地图·图表·告警", RGBColor(0xAA, 0x77, 0xFF)),
    ]

    card_w = Inches(2.15)
    gap = Inches(0.25)
    total_w = len(flow_steps) * card_w + (len(flow_steps) - 1) * gap
    x_start = (SLIDE_WIDTH - total_w) / 2

    for i, (text, color) in enumerate(flow_steps):
        x = x_start + i * (card_w + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), card_w, Inches(1.4),
                  fill_color=CARD_BG, line_color=color, line_width=Pt(2))
        add_textbox(slide, x + Inches(0.1), Inches(1.95), card_w - Inches(0.2), Inches(1.1),
                    text, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        if i < len(flow_steps) - 1:
            arrow_x = x + card_w + Inches(0.02)
            add_textbox(slide, arrow_x, Inches(2.2), gap - Inches(0.04), Inches(0.5),
                        "→", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # 技术栈
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(0.5),
                "核心技术栈", font_size=22, color=WHITE, bold=True)

    tech_cards = [
        ("后端框架", "Python + Flask", ACCENT_BLUE),
        ("数据库", "MongoDB", ACCENT_CYAN),
        ("前端展示", "ECharts + Leaflet", GOLD),
        ("AI / LLM", "多提供商支持", RGBColor(0xAA, 0x77, 0xFF)),
        ("浏览器自动化", "Playwright", RGBColor(0x55, 0xBB, 0x55)),
        ("即时通讯", "Telegram 监控", RED_ALERT),
    ]

    card_w = Inches(1.88)
    gap = Inches(0.2)
    x_start = Inches(0.8)
    for i, (label, value, color) in enumerate(tech_cards):
        x = x_start + i * (card_w + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.1), card_w, Inches(1.2),
                  fill_color=CARD_BG, line_color=color, line_width=Pt(1))
        add_textbox(slide, x + Inches(0.1), Inches(4.2), card_w - Inches(0.2), Inches(0.4),
                    label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(4.55), card_w - Inches(0.2), Inches(0.5),
                    value, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # 系统规模
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.9), Inches(1.2),
              fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(1))
    add_textbox(slide, Inches(1.1), Inches(5.75), Inches(11.3), Inches(0.4),
                "系统规模", font_size=18, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(1.1), Inches(6.2), Inches(11.3), Inches(0.4),
                "160+ REST API 端点   |   4000+ 行后端 API 代码   |   94KB 主仪表盘页面   |   40KB+ 站点解析器",
                font_size=14, color=LIGHT_GRAY)

    add_bottom_bar(slide)
    add_page_number(slide, 10, total_slides)

    # ====================================================
    # 第11页：未来升级 + 成本
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                "三、未来升级  &  四、建设成本", font_size=36, color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(3.5), Inches(0.04), fill_color=ACCENT_BLUE)

    # 上半部分 - 三个升级方向
    upgrades = [
        ("短视频语音分析", ACCENT_BLUE, '让系统"听懂"境外短视频\n视频语音自动转文字分析'),
        ("企业微信推送", ACCENT_CYAN, "高风险事件实时推送\n直达领导和民警手机"),
        ("人员关系图谱", GOLD, "涉案人员关系网络绘制\n多维度关联分析"),
    ]

    card_w = Inches(3.7)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (title, color, desc) in enumerate(upgrades):
        x = x_start + i * (card_w + gap)
        add_decorated_card(slide, x, Inches(1.5), card_w, Inches(1.8), color)
        add_textbox(slide, x + Inches(0.3), Inches(1.72), card_w - Inches(0.6), Inches(0.4),
                    title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.3), Inches(2.2), card_w - Inches(0.6), Inches(0.9),
                    desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # 下半部分 - 三个成本指标
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.03), fill_color=ACCENT_BLUE)

    metrics = [
        ("软件授权费", "¥0", "自主研发，无需购买\n商业软件授权", ACCENT_CYAN),
        ("月运行成本", "< ¥1,000", "开发阶段实测\n大几百元/月", GOLD),
        ("二期升级周期", "2-3 个月", "核心功能已基本成型\n需申请脱产专项", ACCENT_BLUE),
    ]

    for i, (label, value, desc, color) in enumerate(metrics):
        x = x_start + i * (card_w + gap)
        add_decorated_card(slide, x, Inches(3.9), card_w, Inches(2.8), color)
        add_textbox(slide, x + Inches(0.2), Inches(4.1), card_w - Inches(0.4), Inches(0.35),
                    label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(4.5), card_w - Inches(0.4), Inches(0.7),
                    value, font_size=38, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(5.35), card_w - Inches(0.4), Inches(0.9),
                    desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide)
    add_page_number(slide, 11, total_slides)

    # ====================================================
    # 第12页：总结
    # ====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.05), fill_color=ACCENT_BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(0.06), Inches(4.0), fill_color=ACCENT_BLUE)

    add_textbox(slide, Inches(2.0), Inches(1.5), Inches(9), Inches(0.8),
                "总结", font_size=40, color=WHITE, bold=True)

    summary_items = [
        ("花钱少", "软件零成本，月运行仅几百元"),
        ("见效快", "核心功能已成型，即刻投入使用"),
        ("极其安全", "只进不出，数据自主可控"),
        ("固化经验", "先进搜集经验流程化、系统化"),
        ("解放警力", "机器搜集，人做研判，多出战果"),
    ]

    y = Inches(2.5)
    for title, desc in summary_items:
        add_textbox(slide, Inches(2.0), y, Inches(2.5), Inches(0.5),
                    f"✓  {title}", font_size=22, color=ACCENT_CYAN, bold=True)
        add_textbox(slide, Inches(4.5), y + Inches(0.03), Inches(7), Inches(0.5),
                    f"—  {desc}", font_size=18, color=LIGHT_GRAY)
        y += Inches(0.65)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.8), Inches(9.5), Inches(0.8),
              fill_color=SECTION_BLUE, line_color=ACCENT_CYAN, line_width=Pt(2))
    add_textbox(slide, Inches(2.0), Inches(5.85), Inches(9.5), Inches(0.7),
                '皇岗站迈向"亚洲一流" · 实现智慧管控的强力抓手',
                font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.5),
                "谢谢各位领导！", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide)
    add_page_number(slide, 12, total_slides)

    # ====================================================
    # 保存
    # ====================================================
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "全球舆情态势感知平台-讲解PPT.pptx")
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    create_ppt()
